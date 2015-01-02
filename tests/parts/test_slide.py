# encoding: utf-8

"""
Test suite for pptx.parts.slide module
"""

from __future__ import absolute_import

import pytest

from pptx.chart.data import ChartData
from pptx.enum.base import EnumValue
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.opc.package import Part
from pptx.oxml.parts.slide import CT_Slide
from pptx.oxml.shapes.autoshape import CT_Shape
from pptx.oxml.shapes.groupshape import CT_GroupShape
from pptx.package import Package
from pptx.parts.chart import ChartPart
from pptx.parts.image import ImagePart
from pptx.parts.slide import BaseSlide, Slide, _SlidePlaceholders
from pptx.parts.slidelayout import SlideLayout
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.shapes.shapetree import SlideShapeTree

from ..unitutil.cxml import element
from ..unitutil.file import absjoin, test_file_dir

from ..unitutil.mock import (
    class_mock, function_mock, initializer_mock, instance_mock, method_mock,
    property_mock
)


class DescribeBaseSlide(object):

    def it_provides_access_to_its_spTree_element_to_help_ShapeTree(
            self, slide):
        spTree = slide.spTree
        assert isinstance(spTree, CT_GroupShape)

    def it_knows_its_name(self, name_fixture):
        base_slide, expected_value = name_fixture
        assert base_slide.name == expected_value

    def it_can_add_an_image_part(self, image_part_fixture):
        slide, image_file, image_part_, rId_ = image_part_fixture

        image_part, rId = slide.get_or_add_image_part(image_file)

        slide._package.get_or_add_image_part.assert_called_once_with(
            image_file
        )
        slide.relate_to.assert_called_once_with(image_part_, RT.IMAGE)
        assert image_part is image_part_
        assert rId is rId_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def image_part_fixture(
            self, partname_, package_, image_part_, relate_to_):
        slide = BaseSlide(partname_, None, None, package_)
        image_file, rId = 'foobar.png', 'rId6'
        package_.get_or_add_image_part.return_value = image_part_
        relate_to_.return_value = rId
        return slide, image_file, image_part_, rId

    @pytest.fixture
    def name_fixture(self):
        sld_cxml, expected_value = 'p:sld/p:cSld{name=Foobar}', 'Foobar'
        sld = element(sld_cxml)
        base_slide = BaseSlide(None, None, sld, None)
        return base_slide, expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def image_part_(self, request):
        return instance_mock(request, ImagePart)

    @pytest.fixture
    def package_(self, request):
        return instance_mock(request, Package)

    @pytest.fixture
    def partname_(self):
        return PackURI('/foo/bar.xml')

    @pytest.fixture
    def relate_to_(self, request):
        return method_mock(request, BaseSlide, 'relate_to')

    @pytest.fixture
    def slide(self):
        sld = element('p:sld/p:cSld/p:spTree')
        return BaseSlide(None, None, sld, None)


class DescribeSlide(object):

    def it_can_add_a_chart_part(self, add_chart_part_fixture):
        slide, chart_type_, chart_data_ = add_chart_part_fixture[:3]
        ChartPart_, chart_part_, package_, rId = add_chart_part_fixture[3:]

        _rId = slide.add_chart_part(chart_type_, chart_data_)

        ChartPart_.new.assert_called_once_with(
            chart_type_, chart_data_, package_
        )
        slide.relate_to.assert_called_once_with(slide, chart_part_, RT.CHART)
        assert _rId is rId

    def it_provides_access_to_the_shapes_on_the_slide(self, shapes_fixture):
        slide, SlideShapeTree_, slide_shape_tree_ = shapes_fixture
        shapes = slide.shapes
        SlideShapeTree_.assert_called_once_with(slide)
        assert shapes is slide_shape_tree_

    def it_provides_access_to_its_placeholders(self, placeholders_fixture):
        slide, _SlidePlaceholders_, slide_placeholders_ = (
            placeholders_fixture
        )
        placeholders = slide.placeholders
        _SlidePlaceholders_.assert_called_once_with(slide)
        assert placeholders is slide_placeholders_

    def it_can_create_a_new_slide(self, new_fixture):
        slide_layout_, partname_, package_ = new_fixture[:3]
        Slide_init_, slide_elm_, shapes_, relate_to_ = new_fixture[3:]

        slide = Slide.new(slide_layout_, partname_, package_)

        Slide_init_.assert_called_once_with(
            partname_, CT.PML_SLIDE, slide_elm_, package_
        )
        shapes_.clone_layout_placeholders.assert_called_once_with(
            slide_layout_
        )
        relate_to_.assert_called_once_with(
            slide, slide_layout_, RT.SLIDE_LAYOUT
        )
        assert isinstance(slide, Slide)

    def it_knows_the_slide_layout_it_inherits_from(self, layout_fixture):
        slide, slide_layout_ = layout_fixture
        slide_layout = slide.slide_layout
        slide.part_related_by.assert_called_once_with(RT.SLIDE_LAYOUT)
        assert slide_layout is slide_layout_

    def it_knows_the_minimal_element_xml_for_a_slide(self, slide):
        path = absjoin(test_file_dir, 'minimal_slide.xml')
        sld = CT_Slide.new()
        with open(path, 'r') as f:
            expected_xml = f.read()
        assert sld.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def add_chart_part_fixture(
            self, package_, chart_type_, chart_data_, ChartPart_,
            chart_part_, rId, relate_to_):
        slide = Slide(None, None, None, package_)
        return (
            slide, chart_type_, chart_data_, ChartPart_, chart_part_,
            package_, rId
        )

    @pytest.fixture
    def layout_fixture(self, slide_layout_, part_related_by_):
        slide = Slide(None, None, None, None)
        return slide, slide_layout_

    @pytest.fixture
    def placeholders_fixture(
            self, _SlidePlaceholders_, slide_placeholders_):
        slide = Slide(None, None, None, None)
        return slide, _SlidePlaceholders_, slide_placeholders_

    @pytest.fixture
    def shapes_fixture(self, SlideShapeTree_, slide_shape_tree_):
        slide = Slide(None, None, None, None)
        return slide, SlideShapeTree_, slide_shape_tree_

    @pytest.fixture
    def new_fixture(
            self, slide_layout_, partname_, package_, Slide_init_,
            CT_Slide_, slide_elm_, shapes_prop_, shapes_,
            relate_to_):
        return (
            slide_layout_, partname_, package_, Slide_init_, slide_elm_,
            shapes_, relate_to_
        )

    # fixture components -----------------------------------

    @pytest.fixture
    def ChartPart_(self, request, chart_part_):
        ChartPart_ = class_mock(request, 'pptx.parts.slide.ChartPart')
        ChartPart_.new.return_value = chart_part_
        return ChartPart_

    @pytest.fixture
    def chart_data_(self, request):
        return instance_mock(request, ChartData)

    @pytest.fixture
    def chart_part_(self, request):
        return instance_mock(request, ChartPart)

    @pytest.fixture
    def chart_type_(self, request):
        return instance_mock(request, EnumValue)

    @pytest.fixture
    def CT_Slide_(self, request, slide_elm_):
        CT_Slide_ = class_mock(request, 'pptx.parts.slide.CT_Slide')
        CT_Slide_.new.return_value = slide_elm_
        return CT_Slide_

    @pytest.fixture
    def package_(self, request):
        return instance_mock(request, Package)

    @pytest.fixture
    def part_related_by_(self, request, slide_layout_):
        return method_mock(
            request, Slide, 'part_related_by',
            return_value=slide_layout_
        )

    @pytest.fixture
    def partname_(self, request):
        return instance_mock(request, PackURI)

    @pytest.fixture
    def relate_to_(self, request, rId):
        return method_mock(
            request, Part, 'relate_to', autospec=True, return_value=rId
        )

    @pytest.fixture
    def rId(self):
        return 'rId42'

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, SlideShapeTree)

    @pytest.fixture
    def shapes_prop_(self, request, shapes_):
        return property_mock(request, Slide, 'shapes', return_value=shapes_)

    @pytest.fixture
    def slide(self):
        return Slide(None, None, None, None)

    @pytest.fixture
    def _SlidePlaceholders_(self, request, slide_placeholders_):
        return class_mock(
            request, 'pptx.parts.slide._SlidePlaceholders',
            return_value=slide_placeholders_
        )

    @pytest.fixture
    def SlideShapeTree_(self, request, slide_shape_tree_):
        return class_mock(
            request, 'pptx.parts.slide.SlideShapeTree',
            return_value=slide_shape_tree_
        )

    @pytest.fixture
    def slide_elm_(self, request):
        return instance_mock(request, CT_Slide)

    @pytest.fixture
    def Slide_init_(self, request):
        return initializer_mock(request, Slide)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slide_placeholders_(self, request):
        return instance_mock(request, _SlidePlaceholders)

    @pytest.fixture
    def slide_shape_tree_(self, request):
        return instance_mock(request, SlideShapeTree)


class Describe_SlidePlaceholders(object):

    def it_constructs_a_slide_placeholder_for_a_placeholder_shape(
            self, factory_fixture):
        slide_placeholders, ph_elm_ = factory_fixture[:2]
        SlideShapeFactory_, slide_placeholder_ = factory_fixture[2:]
        slide_placeholder = slide_placeholders._shape_factory(ph_elm_)
        SlideShapeFactory_.assert_called_once_with(
            ph_elm_, slide_placeholders
        )
        assert slide_placeholder is slide_placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def factory_fixture(
            self, ph_elm_, SlideShapeFactory_, slide_placeholder_):
        slide_placeholders = _SlidePlaceholders(None)
        return (
            slide_placeholders, ph_elm_, SlideShapeFactory_,
            slide_placeholder_
        )

    # fixture components ---------------------------------------------

    @pytest.fixture
    def slide_placeholder_(self, request):
        return instance_mock(request, SlidePlaceholder)

    @pytest.fixture
    def SlideShapeFactory_(self, request, slide_placeholder_):
        return function_mock(
            request, 'pptx.parts.slide.SlideShapeFactory',
            return_value=slide_placeholder_
        )

    @pytest.fixture
    def ph_elm_(self, request):
        return instance_mock(request, CT_Shape)
